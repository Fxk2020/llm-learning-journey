#!/usr/bin/env python3
"""
将 Word 文档转换为 PPT，使用指定的 PPT 模板
答案放在备注中，不在页面展示
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from docx import Document
import re
import copy

# 文件路径
TEMPLATE_PATH = "2025企业PPT模板-内页版本.pptx"
DOCX_PATH = "2026江苏事业有成-套题刷题-12小时-判断【教师版】.docx"
OUTPUT_PATH = "2026江苏事业有成-判断题PPT.pptx"

def extract_questions_from_docx(docx_path):
    """从 Word 文档中提取题目内容（支持选择题和判断题）"""
    doc = Document(docx_path)
    questions = []
    
    # 获取所有段落文本
    all_text = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            all_text.append(text)
    
    print(f"文档共有 {len(all_text)} 个段落")
    
    current_question = None
    collecting_options = False
    
    i = 0
    while i < len(all_text):
        text = all_text[i]
        
        # 检测题目开始 - 数字开头
        match = re.match(r'^(\d+)[\.、．。]\s*(.+)', text)
        if match:
            # 保存上一题
            if current_question:
                questions.append(current_question)
            
            current_question = {
                'number': match.group(1),
                'question': match.group(2),
                'options': [],
                'answer': None,
                'explanation': ''
            }
            collecting_options = True
            i += 1
            continue
        
        if current_question:
            # 检测选项行 (A. B. C. D. 或 A、B、C、D)
            option_match = re.match(r'^([A-D])[\.、．。\s]\s*(.+)', text)
            if option_match:
                current_question['options'].append(f"{option_match.group(1)}. {option_match.group(2)}")
                i += 1
                continue
            
            # 检测答案行
            answer_match = re.match(r'【答案】\s*([A-D对错√×])', text)
            if answer_match:
                current_question['answer'] = answer_match.group(1)
                collecting_options = False
                i += 1
                continue
            
            # 检测解析行
            exp_match = re.match(r'【解析】\s*(.+)', text)
            if exp_match:
                current_question['explanation'] = exp_match.group(1)
                i += 1
                continue
            
            # 检测其他标记行（跳过）
            if text.startswith('【') and '】' in text:
                i += 1
                continue
            
            # 如果还在收集选项/题目内容
            if collecting_options and not text.startswith('【'):
                # 可能是题目的续行（带①②③④⑤的行）
                if re.match(r'^[①②③④⑤⑥⑦⑧⑨⑩]', text):
                    current_question['question'] += '\n' + text
                elif current_question['options']:
                    # 选项的续行
                    current_question['options'][-1] += text
                else:
                    # 题目续行
                    current_question['question'] += '\n' + text
        
        i += 1
    
    # 保存最后一题
    if current_question:
        questions.append(current_question)
    
    return questions

def create_ppt_from_template(template_path, questions, output_path):
    """使用模板创建新的 PPT - 答案放在备注中"""
    # 加载模板
    prs = Presentation(template_path)
    
    # 获取幻灯片尺寸
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    
    # 找一个适合内容页的布局 (选择 Blank 布局)
    content_layout = None
    for i, layout in enumerate(prs.slide_layouts):
        if 'Blank' in layout.name or '空白' in layout.name:
            content_layout = layout
            break
    
    if content_layout is None:
        content_layout = prs.slide_layouts[17]
    
    # 删除所有现有幻灯片
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]
    
    # 配置样式 - 更美观的边距和位置
    margin_top = Inches(0.8)
    margin_left = Inches(1.0)
    margin_right = Inches(1.0)
    content_width = slide_width - margin_left - margin_right
    
    # 主题色
    question_color = RGBColor(45, 45, 45)  # 深灰色
    option_color = RGBColor(60, 60, 60)  # 灰色
    
    for q in questions:
        # 添加新幻灯片
        slide = prs.slides.add_slide(content_layout)
        
        # ========== 题目内容 ==========
        question_box = slide.shapes.add_textbox(margin_left, Inches(0.4), content_width, Inches(2.5))
        question_tf = question_box.text_frame
        question_tf.word_wrap = True
        
        q_p = question_tf.paragraphs[0]
        q_p.text = q['question']
        q_p.font.size = Pt(20)
        q_p.font.name = "微软雅黑"
        q_p.font.color.rgb = question_color
        q_p.line_spacing = 1.5
        
        # ========== 选项 ==========
        if q['options']:
            options_top = Inches(0.4) + Inches(2.7)
            options_box = slide.shapes.add_textbox(margin_left + Inches(0.3), options_top, content_width - Inches(0.3), Inches(3.5))
            options_tf = options_box.text_frame
            options_tf.word_wrap = True
            
            for idx, opt in enumerate(q['options']):
                if idx == 0:
                    opt_p = options_tf.paragraphs[0]
                else:
                    opt_p = options_tf.add_paragraph()
                
                opt_p.text = opt
                opt_p.font.size = Pt(20)
                opt_p.font.name = "微软雅黑"
                opt_p.font.color.rgb = option_color
                opt_p.line_spacing = 1.6
                opt_p.space_before = Pt(12)
        
        # ========== 答案和解析放入备注 ==========
        notes_slide = slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        
        notes_content = []
        if q['answer']:
            notes_content.append(f"【答案】{q['answer']}")
        if q.get('explanation'):
            notes_content.append(f"\n【解析】{q['explanation']}")
        
        if notes_content:
            notes_tf.text = '\n'.join(notes_content)
    
    # 保存 PPT
    prs.save(output_path)
    print(f"\nPPT 已保存到: {output_path}")
    print(f"共生成 {len(prs.slides)} 张幻灯片，包含 {len(questions)} 道题目")
    print("答案和解析已放入幻灯片备注中")

if __name__ == "__main__":
    print("=" * 60)
    print("提取 Word 文档内容")
    print("=" * 60)
    questions = extract_questions_from_docx(DOCX_PATH)
    print(f"\n共提取 {len(questions)} 道题目")
    
    # 统计有答案的题目
    with_answer = sum(1 for q in questions if q['answer'])
    with_options = sum(1 for q in questions if q['options'])
    print(f"有答案的题目: {with_answer} 道")
    print(f"有选项的题目: {with_options} 道")
    
    if questions:
        print(f"\n前3题示例:")
        for q in questions[:3]:
            print(f"\n  {q['number']}. {q['question'][:60]}...")
            if q['options']:
                for opt in q['options'][:2]:
                    print(f"     {opt[:50]}...")
            print(f"     答案: {q['answer']} (将放入备注)")
    
    print("\n" + "=" * 60)
    print("生成 PPT")
    print("=" * 60)
    create_ppt_from_template(TEMPLATE_PATH, questions, OUTPUT_PATH)
